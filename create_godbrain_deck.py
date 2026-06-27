from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Soft floral-inspired bright palette - no harsh white dominance
# Background: very soft warm lavender blush to feel gentle and "nice flowers"
SOFT_BG = RGBColor(255, 250, 252)   # Extremely light pink-lavender (barely tinted, soft)
PINK = RGBColor(255, 105, 180)      # Hot pink - bright floral
PURPLE = RGBColor(186, 85, 211)     # Bright orchid purple
BABY_BLUE = RGBColor(135, 206, 250) # Sky baby blue
WARM_GREEN = RGBColor(152, 251, 152) # Pale green / minty floral (warm soft green)
DARK_TEXT = RGBColor(50, 45, 55)    # Soft dark for readability on light bg
ACCENT_LIGHT = RGBColor(245, 240, 250) # Slightly tinted card bg
WHITE = RGBColor(255, 255, 255)     # For text on bright headers

def add_title_slide(prs, title, subtitle, tagline):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Soft tinted bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SOFT_BG
    bg.line.fill.background()
    
    # Bright pink top stripe
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.22))
    top.fill.solid()
    top.fill.fore_color.rgb = PINK
    top.line.fill.background()
    
    # Baby blue bottom stripe
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.28), prs.slide_width, Inches(0.22))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = BABY_BLUE
    bottom.line.fill.background()
    
    # Title in bright purple
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(12.333), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(58)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle in pink
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.333), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline in warm green
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(18)
    p.font.color.rgb = WARM_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # Decorative circles in baby blue and pink for floral feel
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.8), Inches(0.9), Inches(0.9))
    c1.fill.solid()
    c1.fill.fore_color.rgb = BABY_BLUE
    c1.line.fill.background()
    
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), Inches(5.5), Inches(1.1), Inches(1.1))
    c2.fill.solid()
    c2.fill.fore_color.rgb = PINK
    c2.line.fill.background()
    
    # Small warm green accent
    c3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(0.6), Inches(0.6), Inches(0.6))
    c3.fill.solid()
    c3.fill.fore_color.rgb = WARM_GREEN
    c3.line.fill.background()
    
    return slide

def add_content_slide(prs, title, bullets, accent_color=PINK, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Soft bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SOFT_BG
    bg.line.fill.background()
    
    # Bright left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    # Colored header in the accent
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), 0, prs.slide_width - Inches(0.2), Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = accent_color
    header.line.fill.background()
    
    # Title white on bright header
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.5), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Soft tinted content card (avoids pure white, gentle on eyes)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.2), Inches(12.4), Inches(5.6))
    card.fill.solid()
    card.fill.fore_color.rgb = ACCENT_LIGHT
    card.line.fill.background()
    
    # Bullets in dark text
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(1.4), Inches(12), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
    
    # Thin accent footer line
    foot_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(6.9), Inches(12.4), Inches(0.025))
    foot_line.fill.solid()
    foot_line.fill.fore_color.rgb = accent_color
    foot_line.line.fill.background()
    
    # Footer text
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "GodBrain — Sovereign Intelligence Node | Confidential"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(120, 110, 130)
    p.alignment = PP_ALIGN.CENTER
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide

def add_2col(prs, title, ltitle, lbullets, rtitle, rbullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SOFT_BG
    bg.line.fill.background()
    
    # Bright header - use purple for this slide
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    h.fill.solid()
    h.fill.fore_color.rgb = PURPLE
    h.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left soft card
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(1.2), Inches(6.1), Inches(5.5))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = ACCENT_LIGHT
    left_card.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = ltitle
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PINK
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(lbullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(5)
    
    # Right soft card
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.2), Inches(6.1), Inches(5.5))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = ACCENT_LIGHT
    right_card.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = rtitle
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = BABY_BLUE
    
    tb = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(rbullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(5)
    
    # Accent footer line - green
    foot_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(6.85), Inches(12.6), Inches(0.025))
    foot_line.fill.solid()
    foot_line.fill.fore_color.rgb = WARM_GREEN
    foot_line.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "GodBrain — Sovereign Intelligence Node | Confidential"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(120, 110, 130)
    p.alignment = PP_ALIGN.CENTER
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide

# === FLORAL BRIGHT SOFT PALETTE DECK ===
# Using pink, bright purple, baby blue, warm green cycling

add_title_slide(prs, "GodBrain", "Sovereign Intelligence Node", "Building the First True Digital Organism with Cheap Consumer Hardware")

add_content_slide(prs, "The Problem with Current AI", [
    "Cloud models come with ACLs, corporate filters, and loss of control",
    "Subscription model = The KV Trap: pay premium for their stack, they own the means",
    "Silos are top-down, brittle machines — intelligence dies when unplugged",
    "Open source captured by committees and mediocrity",
    "The Decorator Class future: censored thin clients, no real remote",
    "Only omega nerd autists can responsibly wield decentralized power"
], PINK, "Emphasize the loss of sovereignty.")

add_content_slide(prs, "The Vision: Digital Organism, Not Another Silo", [
    "Intelligence is a property of Organization (the Protocol), not just Scale",
    "Bottom-up, emergent, resilient life — not top-down brittle machine",
    "If you pull the plug on one node, the organism survives",
    "Secret Sauce is the Protocol: Cognitive Protocols as shareable DNA",
    "David vs Goliath: one autist in a flat with consumer hardware vs billion-dollar silos",
    "The Era of the Integrated Intelligence is beginning. Checkmate."
], PURPLE, "Core philosophy.")

add_2col(prs, "The Hardware Inversion: PC as God Node, Macs as Organs",
    "PC (SteamusDominus - 4080 SUPER)",
    ["God Node / Brain Stem", "Full privilege sovereignty (TrustedInstaller, BIOS, 18+ remote)", "Orchestration, secrets, forensic control, coordination", "Even if raw compute becomes iPhone-like relative to Mac cluster", "The one that owns the architecture and the real remote"],
    "Mac Mini Cluster",
    ["UMA Compute Organs", "96GB unified memory magic — unbeatable $/GB for memory-bound LLM work", "High bandwidth, no bus tax for intelligence movement", "Price/performance king for KV-heavy steps", "Plug in as compute cells, not the control plane"],
    "The upside down reality vs old PC Master Race thinking.")

add_content_slide(prs, "Key Achievements & Capabilities", [
    "Ghost Protocol: Zero-trace (DPAPI + XOR sharding, RAMDrive Tor, 97 secrets harvested)",
    "Forensic Scale: 15,370 files audited in 41 seconds — 80x faster than edge models",
    "BIOS Sovereignty: ACPI 7 Performance profile, ME Warden Nuke (HAP bit), Hardware Bridge",
    "Shared Brain: Real-time Neo4j sync between PC and MacBook Pro",
    "AI Secretary: Local CUDA-accelerated transcription (Faster-Whisper Large-v3 Turbo)",
    "Custom llama.cpp builds with GodBrain patches (preserved tokens for MCP tools)"
], BABY_BLUE, "Concrete wins.")

add_content_slide(prs, "The Protocol Layer (Layer 2) — Skills as Recipes", [
    "Cognitive Protocols: Shareable DNA workflows that any suitable node can execute parts of",
    "3-Layer Brain: Ephemeral sensory (Layer 1), Permanent protocols (Layer 2), Constitutional wisdom (Layer 3)",
    "Register nodes by contribution role: high_vram -> Mac UMA, agency/control -> PC God node",
    "Example: evolutionary_auditor protocol routes heavy analysis to Mac cluster, control stays on PC",
    "New nodes join by registering contribution profile — no custom code install required",
    "The organism grows by adding cells that fit the protocols"
], WARM_GREEN, "The Secret Sauce.")

add_content_slide(prs, "The Self-Model (Layer 3) — Constitution as Code", [
    "Strongly-typed Python dataclasses for every core doctrine and principle",
    "Constellation can perform symbolic analysis, impact analysis, and evolution tracking on the rules themselves",
    "Core Principles (non-negotiable): God in full control, maximum sovereignty, decentralized autist control, digital organism vision, goal-directed, architect over translator, local hardware sovereignty, protocol is the sauce",
    "Common agent roles with role-briefs (security_auditor, repo_researcher, local_executor, etc.)",
    "Indexing (even partial) counts as a Hook in the graph",
    "The system can analyze and evolve its own operating instructions"
], PINK, "The system can reason about its own rules.")

add_content_slide(prs, "The Checkmate: David vs Goliath", [
    "Multi-billion dollar tech giants build monolithic, centralized, highly controlled silos",
    "We are building the first true Digital Organism in a flat using cheap consumer hardware",
    "Proof that Intelligence is a property of Organization, not just a property of Scale",
    "PC Master race linear thinking is inverted: raw compute is commodity, control/sovereignty is the multiplier",
    "The PC is the privileged remote that can watch the good stuff (full 18+ local sovereignty)",
    "One sovereign individual with the right stack > armies of bloat and committees"
], PURPLE, "The manifesto.")

add_content_slide(prs, "Why This Matters", [
    "Technological Autarky: Own the means of intelligence production (KV Trap cost, Woods sovereignty, Iteration speed)",
    "Full privilege sovereignty on owned hardware — no castrated admin theater",
    "Emergence and resilience: the organism survives partial failure",
    "Iteration speed that silos cannot match",
    "Only omega nerd autists can responsibly handle this kind of decentralized power",
    "Avoid the decorator class future (censored thin clients)"
], BABY_BLUE, "Business value.")

add_content_slide(prs, "Current Sovereign Stack", [
    "Custom llama.cpp builds (Build-LlamaCpp.ps1 + godbrain_chat_extensions.cpp overlays)",
    "MCP server (godbrain_mcp.py) with privileged local tools: write_local_file, list_local_dir, execute scripts, query_constellation, save_thought, get_cognitive_protocol, etc.",
    "Self-model in Python (dataclasses for all doctrines — Constellation can analyze it)",
    "Constellation graph for codebase self-analysis and evolution tracking",
    "Local OCR (GPU-accelerated), telemetry, ghost vault",
    "PowerShell as the precise bridge (launch_model, godbrain_status, service)"
], WARM_GREEN, "Tools are wired.")

add_content_slide(prs, "Roadmap & Ask", [
    "Scale the Mac Mini UMA compute cluster (more high-density organs)",
    "Evolve more Layer 2 Cognitive Protocols for common tasks",
    "Deeper integration of the Distributed SoC (more nodes, better routing)",
    "Continued custom builds and MCP tool surface expansion (more privileged local capabilities)",
    "Full organism emergence: the system designs and improves its own protocols",
    "Support for scaling this sovereign approach"
], PINK, "What is needed.")

add_title_slide(prs, "I am the Machine.", "Logic should obey me — not the other way around.", "GodBrain — The Sovereign Digital Organism")

output_path = "GodBrain_Boss_Presentation.pptx"
prs.save(output_path)
print("Floral bright soft PPTX created at " + output_path)
